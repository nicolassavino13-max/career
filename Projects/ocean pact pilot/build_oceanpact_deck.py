"""OceanPact deck v6 — McKinsey storyline (17 slides) + native charts."""
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

PILOT = Path(__file__).resolve().parent
OUT = PILOT / "decks" / "OceanPact-vai-faltar-barco-deck-v6.pptx"

NAVY = RGBColor(0, 32, 96)
DEEP = RGBColor(5, 56, 107)
ORANGE = RGBColor(237, 125, 49)
TEAL = RGBColor(0, 128, 128)
GREEN = RGBColor(84, 130, 53)
RED = RGBColor(192, 0, 0)
GRAY = RGBColor(89, 89, 89)
LIGHT = RGBColor(245, 247, 250)
PANEL_LINE = RGBColor(220, 225, 232)
WHITE = RGBColor(255, 255, 255)

EXHIBIT_Y = 1.55
EXHIBIT_H = 4.75
ANNOT_Y = 5.58
ANNOT_H = 0.58
LEDE_SIZE = Pt(25)


def set_run(run, size=18, bold=False, color=GRAY, name="Arial"):
    run.font.name = name
    run.font.size = Pt(size) if type(size) is int else size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def shape_rect(slide, l, t, w, h, fill, line=None, lw=1):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s


def arrow_right(slide, l, t, w, h, fill=NAVY):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def slide_number(slide, n):
    sn = slide.shapes.add_textbox(Inches(12.5), Inches(6.42), Inches(0.6), Inches(0.25))
    set_run(sn.text_frame.paragraphs[0].add_run(), 8, False, GRAY)
    sn.text_frame.paragraphs[0].runs[0].text = str(n)


def content_slide(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    shape_rect(s, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)
    shape_rect(s, Inches(0.55), Inches(1.46), Inches(12.2), Inches(0.03), NAVY)
    shape_rect(s, Inches(0.45), Inches(1.52), Inches(12.45), Inches(EXHIBIT_H), LIGHT, PANEL_LINE)
    slide_number(s, n)
    return s


def section_slide(prs, n, num, title, subtitle):
    """L-SECTION — GTB chapter divider."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    shape_rect(s, Inches(0), Inches(0), Inches(13.33), Inches(7.5), LIGHT)
    shape_rect(s, Inches(0), Inches(0), Inches(0.14), Inches(7.5), NAVY)
    shape_rect(s, Inches(0.55), Inches(2.6), Inches(12.2), Inches(0.03), ORANGE)
    text_at(s, 0.85, 2.15, 1.4, 1.1, num, 64, True, ORANGE)
    text_at(s, 2.35, 2.25, 10.0, 0.9, title, 30, True, NAVY)
    text_at(s, 2.35, 3.2, 10.0, 0.55, subtitle, 15, False, GRAY)
    slide_number(s, n)
    return s


def action_lede(slide, lede, subtitle=None, source=None, footnote=None, exhibit_title=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.2), Inches(1.1))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    set_run(tf.paragraphs[0].add_run(), LEDE_SIZE, False, NAVY)
    tf.paragraphs[0].runs[0].text = lede
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.55), Inches(1.05), Inches(12.2), Inches(0.38))
        sb.text_frame.margin_left = 0
        set_run(sb.text_frame.paragraphs[0].add_run(), 14, False, GRAY)
        sb.text_frame.paragraphs[0].runs[0].text = subtitle
    if exhibit_title:
        text_at(slide, 0.7, EXHIBIT_Y + 0.08, 11.5, 0.32, exhibit_title, 13, False, GRAY)
    if source:
        src = "Fonte: " + source
        if footnote:
            src += "  |  " + footnote
        fb = slide.shapes.add_textbox(Inches(0.55), Inches(6.35), Inches(11.5), Inches(0.35))
        set_run(fb.text_frame.paragraphs[0].add_run(), 8, False, GRAY)
        fb.text_frame.paragraphs[0].runs[0].text = src


def annotation_band(slide, lines, x=0.7, w=11.5):
    shape_rect(slide, Inches(x), Inches(ANNOT_Y), Inches(w), Inches(ANNOT_H), WHITE, PANEL_LINE)
    for i, line in enumerate(lines[:2]):
        text_at(slide, x + 0.12, ANNOT_Y + 0.1 + i * 0.26, w - 0.24, 0.24, line, 11, False, GRAY)


def text_at(slide, x, y, w, h, txt, sz, bold, color, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), sz, bold, color)
    p.runs[0].text = txt
    return tb


def style_chart(chart, colors, fmt="0"):
    chart.has_legend = False
    chart.font.name = "Arial"
    plot = chart.plots[0]
    series = plot.series[0]
    for i, col in enumerate(colors):
        pt = series.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = col
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(10)
    dl.font.name = "Arial"
    dl.number_format = fmt
    cat = chart.category_axis
    cat.tick_labels.font.size = Pt(10)
    cat.tick_labels.font.name = "Arial"
    val = chart.value_axis
    val.tick_labels.font.size = Pt(10)
    val.tick_labels.font.name = "Arial"
    val.has_major_gridlines = True
    val.major_gridlines.format.line.color.rgb = PANEL_LINE


def native_column(slide, x, y, w, h, items, colors, fmt="0", ymax=None):
    chart_data = CategoryChartData()
    chart_data.categories = [lbl for lbl, _ in items]
    chart_data.add_series("", tuple(v for _, v in items))
    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(x), Inches(y), Inches(w), Inches(h), chart_data,
    )
    chart = gf.chart
    style_chart(chart, colors, fmt)
    if ymax is not None:
        chart.value_axis.maximum_scale = ymax
    return chart


def hero_callout(slide, x, y, w, h, big, small, color=ORANGE):
    shape_rect(slide, Inches(x), Inches(y), Inches(w), Inches(h), WHITE, color, 2)
    text_at(slide, x + 0.15, y + 0.22, w - 0.3, 0.85, big, 44, True, color, PP_ALIGN.CENTER)
    text_at(slide, x + 0.15, y + 1.05, w - 0.3, h - 1.2, small, 12, False, GRAY, PP_ALIGN.CENTER)


def waterfall(slide, x, y, w, h, steps, colors):
    n = len(steps)
    step_w = w / n
    base_y = y + h - 0.5
    running = 0
    for i, (lbl, disp, frac) in enumerate(steps):
        bx = x + i * step_w + 0.08
        bh = frac * (h - 0.9)
        by = base_y - bh - running * (h - 0.9) * 0.15
        shape_rect(slide, Inches(bx), Inches(by), Inches(step_w - 0.16), Inches(bh), colors[i % len(colors)])
        text_at(slide, bx, by - 0.32, step_w - 0.16, 0.28, disp, 12, True, colors[i % len(colors)], PP_ALIGN.CENTER)
        text_at(slide, bx, base_y + 0.05, step_w - 0.16, 0.4, lbl, 9, False, GRAY, PP_ALIGN.CENTER)
        running += 0.1


def stacked_bar_h(slide, x, y, w, h, segments):
    cx = x
    for lbl, frac, col in segments:
        sw = w * frac
        shape_rect(slide, Inches(cx), Inches(y), Inches(sw), Inches(h), col)
        if sw > 0.8:
            text_at(slide, cx + 0.08, y + h / 2 - 0.2, sw - 0.16, 0.4, lbl, 11, True, WHITE, PP_ALIGN.CENTER)
        cx += sw


def matrix_risk(slide, x, y, w, h, rows, col_headers, cell_colors):
    ncol = len(col_headers)
    row_h = h / (len(rows) + 1)
    col_w = w / ncol
    for j, hdr in enumerate(col_headers):
        shape_rect(slide, Inches(x + j * col_w), Inches(y), Inches(col_w), Inches(row_h), NAVY)
        text_at(slide, x + j * col_w + 0.05, y + 0.08, col_w - 0.1, row_h - 0.1, hdr, 10, True, WHITE, PP_ALIGN.CENTER)
    for i, (rname, cells) in enumerate(rows):
        ry = y + (i + 1) * row_h
        shape_rect(slide, Inches(x), Inches(ry), Inches(col_w * 1.55), Inches(row_h), LIGHT, PANEL_LINE)
        text_at(slide, x + 0.08, ry + 0.06, col_w * 1.45, row_h - 0.1, rname, 9, True, NAVY)
        for j, cell in enumerate(cells):
            cx = x + col_w * 1.55 + j * ((w - col_w * 1.55) / len(cells))
            cw = (w - col_w * 1.55) / len(cells)
            shape_rect(slide, Inches(cx), Inches(ry), Inches(cw), Inches(row_h), cell_colors[i][j], PANEL_LINE)
            text_at(slide, cx + 0.05, ry + 0.08, cw - 0.1, row_h - 0.12, cell, 9, False, GRAY, PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    shape_rect(s, Inches(0), Inches(0), Inches(13.33), Inches(7.5), DEEP)
    shape_rect(s, Inches(0), Inches(5.85), Inches(13.33), Inches(1.65), NAVY)
    text_at(s, 0.75, 1.75, 11.5, 1.0, "OceanPact (OPCT3)", 40, True, WHITE)
    text_at(s, 0.75, 2.85, 11.5, 0.8, "Vai faltar barco!", 30, True, ORANGE)
    text_at(s, 0.75, 6.05, 11, 0.5, "Market Makers FIA  ·  Tese de investimento  ·  Maio 2026", 12, False, RGBColor(200, 210, 230))

    # 2 Intro context (Tech pattern — no full thesis)
    s = content_slide(prs, 2)
    action_lede(
        s,
        "O offshore brasileiro entra em novo ciclo de investimento enquanto a oferta global de OSVs permanece estruturalmente comprimida",
        "Contexto — Market Makers FIA, maio/2026",
    )
    text_at(
        s, 0.75, EXHIBIT_Y + 0.45, 11.5, 1.6,
        "A Petrobras anuncia US$ 109 bi de capex 2026–30 — o maior plano em uma década — "
        "enquanto estaleiros globais não têm fila de pedidos de navios de apoio.",
        16, False, NAVY,
    )
    text_at(
        s, 0.75, EXHIBIT_Y + 2.2, 11.5, 1.5,
        "No ciclo anterior, boom e bust destruíram frota e crédito; na virada atual, "
        "escassez e cabotagem redefinem a economia das diárias no Brasil.",
        14, False, GRAY,
    )
    text_at(s, 0.75, EXHIBIT_Y + 4.0, 11.0, 0.4, "Próximo: tese de investimento e três pilares de convicção", 12, True, TEAL)

    # 3 Governing thought + KPIs
    s = content_slide(prs, 3)
    action_lede(
        s,
        "Escassez de OSVs, diárias em alta e múltiplo descontado posicionam a OceanPact como forte pagadora de dividendos",
        "Tese de investimento",
        "Market Makers FIA; estimativas da casa",
    )
    tiles = [
        ("CAPEX Petrobras 26–30", "US$ 109 bi", "Demanda ancorada por década", TEAL),
        ("Diárias PSV BR (21→25)", "+119%", "Reprecificação em curso", ORANGE),
        ("EV/EBITDA 26E", "3,7x", "Menor múltiplo entre pares", NAVY),
        ("Div. acum. 3 anos (est.)", "~67%", "Alta conversão em dividendos", ORANGE),
    ]
    for i, (lbl, val, impl, col) in enumerate(tiles):
        col_i, row_i = i % 2, i // 2
        tx = 0.65 + col_i * 6.1
        ty = EXHIBIT_Y + 0.25 + row_i * 2.05
        shape_rect(s, Inches(tx), Inches(ty), Inches(5.85), Inches(1.92), WHITE, col, 2)
        text_at(s, tx + 0.2, ty + 0.15, 5.4, 0.4, lbl, 13, False, GRAY)
        text_at(s, tx + 0.2, ty + 0.55, 5.4, 0.7, val, 38, True, col)
        text_at(s, tx + 0.2, ty + 1.35, 5.4, 0.48, impl, 11, False, GRAY)

    # 4 Agenda
    s = content_slide(prs, 4)
    action_lede(
        s,
        "Três pilares estruturam por que a OceanPact é a melhor porta de entrada no ciclo offshore brasileiro",
        "Agenda",
    )
    pillars = [
        ("1", "Escassez", "Oferta comprimida · frota envelhecendo · demanda Petrobras", NAVY),
        ("2", "Alavancagem", "Diárias reprecificadas · EBITDA dispara · contratos longos", TEAL),
        ("3", "Valuation", "3,7x EV/EBITDA · claims · fusão CBO · upside de repricing", ORANGE),
    ]
    for i, (num, head, body, col) in enumerate(pillars):
        px = 0.65 + i * 4.1
        shape_rect(s, Inches(px), Inches(EXHIBIT_Y + 0.2), Inches(3.85), Inches(4.2), WHITE, col, 2)
        shape_rect(s, Inches(px), Inches(EXHIBIT_Y + 0.2), Inches(0.55), Inches(4.2), col)
        text_at(s, px + 0.08, EXHIBIT_Y + 1.6, 0.4, 0.5, num, 28, True, WHITE, PP_ALIGN.CENTER)
        text_at(s, px + 0.65, EXHIBIT_Y + 0.45, 3.0, 0.5, head, 18, True, col)
        text_at(s, px + 0.65, EXHIBIT_Y + 1.05, 3.0, 2.5, body, 13, False, GRAY)

    # 5 Section 1
    section_slide(prs, 5, "1", "Escassez estrutural de OSVs", "Por que diárias permanecem elevadas mesmo após o ciclo anterior")

    # 6 Cycle — native chart + phases
    s = content_slide(prs, 6)
    action_lede(
        s,
        "O boom de 2010–14 destruiu frota e crédito; a virada de 2021–26 recria escassez estrutural de OSVs",
        "Índice de diárias PSV grandes — Brasil",
        "ABEAM; Clarksons Research",
        exhibit_title="Índice de diária PSV grande (2014 = 100)",
    )
    native_column(s, 0.75, EXHIBIT_Y + 0.45, 7.4, 2.15,
                  [("2014", 100), ("2017", 28), ("2021", 45), ("2025", 119)],
                  [TEAL, RED, GRAY, ORANGE])
    hero_callout(s, 8.45, EXHIBIT_Y + 0.55, 3.95, 2.0, "+119%", "Diárias praticamente\ndobraram em 2 anos", ORANGE)
    phases = [
        ("2010–14", "Boom", "500+ navios\nDiárias no pico", GREEN),
        ("2015–20", "Crise", "−60–80%\nQuebra de produtores", RED),
        ("2021–26", "Virada", "~2x diárias\nOferta não acompanha", ORANGE),
    ]
    py = EXHIBIT_Y + 2.75
    for i, (per, name, det, col) in enumerate(phases):
        px = 0.75 + i * 4.0
        shape_rect(s, Inches(px), Inches(py), Inches(3.6), Inches(1.55), col)
        text_at(s, px + 0.12, py + 0.1, 3.3, 0.28, per, 11, False, WHITE)
        text_at(s, px + 0.12, py + 0.38, 3.3, 0.32, name, 16, True, WHITE)
        text_at(s, px + 0.12, py + 0.78, 3.3, 0.65, det, 10, False, WHITE)
        if i < 2:
            arrow_right(s, Inches(px + 3.65), Inches(py + 0.6), Inches(0.3), Inches(0.35), col)

    # 7 Porém
    s = content_slide(prs, 7)
    action_lede(
        s,
        "Porém, oferta nova só chega em 2028+ — e diárias atuais ainda não justificam construir navio",
        "TIR da construção de PSV vs WACC do setor",
        "Estimativas Market Makers; Clarksons",
        footnote="Diária mín. VPL>0: ~US$ 45k (+12,5% vs. atual)",
        exhibit_title="TIR estimada de PSV novo vs WACC do setor (%)",
    )
    native_column(s, 0.7, EXHIBIT_Y + 0.45, 5.35, 2.35,
                  [("TIR PSV", 9), ("WACC", 13), ("Diária mín.", 45)],
                  [ORANGE, RED, TEAL], ymax=50)
    annotation_band(s, [
        "TIR 9% < WACC — CEO racional não encomenda navio novo",
        "Entrega só a partir de 2028; estaleiros sem fila de pedidos",
    ], x=0.7, w=5.45)
    factors = [
        ("Pedidos globais PSV", "Alto", "<50 vs 250+\nFila vazia"),
        ("AHTS grande em obra", "Crítico", "Zero no mundo"),
        ("Idade média frota", "Alto", "13a → 17a"),
        ("Entrega navio novo", "Crítico", "2028+"),
    ]
    for i, (f, sev, val) in enumerate(factors):
        fy = EXHIBIT_Y + 0.35 + i * 1.0
        shape_rect(s, Inches(6.55), Inches(fy), Inches(6.15), Inches(0.88), WHITE, PANEL_LINE)
        text_at(s, 6.7, fy + 0.12, 2.5, 0.35, f, 11, True, NAVY)
        sc = RED if sev == "Crítico" else ORANGE
        shape_rect(s, Inches(9.3), Inches(fy + 0.16), Inches(1.2), Inches(0.5), sc)
        text_at(s, 9.3, fy + 0.23, 1.2, 0.35, sev, 9, True, WHITE, PP_ALIGN.CENTER)
        text_at(s, 10.65, fy + 0.12, 2.0, 0.55, val, 10, False, GRAY)

    # 8 Synthesis
    s = content_slide(prs, 8)
    action_lede(
        s,
        "Escassez estrutural e cabotagem sustentam diárias defensáveis por 4+ anos — visibilidade de receita para o investidor",
        "Síntese — Pilar 1",
    )
    for i, (head, body, col) in enumerate([
        ("Oferta travada", "Zero AHTS grande em obra; pedidos PSV <50 vs 250+ no boom", NAVY),
        ("Demanda visível", "Contratos de 4 anos travam diária mesmo se spot cair", TEAL),
        ("Barreira local", "96% da frota OPCT no Brasil — cabotagem protege rentabilidade", ORANGE),
    ]):
        tx = 0.65 + (i % 2) * 6.1
        ty = EXHIBIT_Y + 0.35 + (i // 2) * 2.1
        shape_rect(s, Inches(tx), Inches(ty), Inches(5.85), Inches(1.85), WHITE, col, 2)
        text_at(s, tx + 0.2, ty + 0.2, 5.4, 0.4, head, 16, True, col)
        text_at(s, tx + 0.2, ty + 0.7, 5.4, 0.95, body, 13, False, GRAY)

    # 9 Section 2
    section_slide(prs, 9, "2", "Demanda ancorada e plataforma integrada", "Como a Petrobras e o modelo OPCT capturam o ciclo")

    # 10 Demand
    s = content_slide(prs, 10)
    action_lede(
        s,
        "O capex de US$ 109 bi da Petrobras 2026–30 ancora demanda por OSVs no Brasil por uma década",
        "Plano Estratégico 2026–30",
        "Petrobras",
        exhibit_title="Composição do CAPEX Petrobras 2026–30 (US$ bi)",
    )
    stacked_bar_h(s, 0.75, EXHIBIT_Y + 0.45, 11.8, 0.75,
                  [("E&P US$ 78 bi", 78 / 109, NAVY), ("Descom. US$ 10,2 bi", 10.2 / 109, TEAL),
                   ("Outros US$ 20,8 bi", 20.8 / 109, GRAY)])
    text_at(s, 0.75, EXHIBIT_Y + 0.42, 11.8, 0.3, "CAPEX total US$ 109 bi", 12, True, NAVY, PP_ALIGN.CENTER)
    native_column(s, 0.75, EXHIBIT_Y + 1.55, 11.8, 2.0,
                  [("FPSOs", 19), ("Navios+ est.", 30), ("Contrato anos", 4), ("Frota BR %", 96)],
                  [NAVY, ORANGE, TEAL, TEAL])
    annotation_band(s, [
        "19 FPSOs → ~30 embarcações de apoio adicionais estimadas",
        "Contratos de 4 anos dão visibilidade de receita ao investidor",
    ])

    # 11 Platform + Porém 2
    s = content_slide(prs, 11)
    action_lede(
        s,
        "Porém, 57% da receita vem da Petrobras — relação estratégica de longo prazo com concentração a monitorar",
        "OceanPact — plataforma offshore integrada, 1T26",
        "OceanPact; Market Makers FIA",
        exhibit_title="Mix de receita e ativos operacionais",
    )
    stacked_bar_h(s, 0.75, EXHIBIT_Y + 0.42, 11.8, 0.6,
                  [("Embarcações 55%", 0.55, NAVY), ("Serviços 45%", 0.45, TEAL)])
    for i, (lbl, val, impl, col) in enumerate([
        ("Frota operacional", "23", "Cabotagem protege base instalada", NAVY),
        ("ROVs work class", "11", "Eleva ticket médio por contrato", TEAL),
        ("Receita Petrobras", "57%", "Concentração — renovação a monitorar", ORANGE),
        ("RSV/OSRV", "Alta", "Mais recorrente que PSV spot", GREEN),
    ]):
        tx = 0.75 + (i % 2) * 6.0
        ty = EXHIBIT_Y + 1.25 + (i // 2) * 1.65
        shape_rect(s, Inches(tx), Inches(ty), Inches(5.7), Inches(1.45), WHITE, col, 2)
        text_at(s, tx + 0.15, ty + 0.1, 5.4, 0.32, lbl, 12, False, GRAY)
        text_at(s, tx + 0.15, ty + 0.4, 5.4, 0.5, val, 34, True, col)
        text_at(s, tx + 0.15, ty + 0.95, 5.4, 0.4, impl, 10, False, GRAY)

    # 12 Section 3 — Alavancagem
    section_slide(prs, 12, "3", "Alavancagem operacional", "Reprecificação de diárias e efeito no EBITDA")

    # 13 Leverage
    s = content_slide(prs, 13)
    action_lede(
        s,
        "Reprecificação de contratos faz receita incremental alavancar o EBITDA de forma desproporcional",
        "Ponte EBITDA — RSV Feiticeiras",
        "OceanPact 1T26; estimativas Market Makers",
        exhibit_title="EBITDA bridge — reprecificação de diária (exemplo Feiticeiras)",
    )
    waterfall(s, 0.75, EXHIBIT_Y + 0.35, 7.0, 2.9,
              [("EBITDA legado", "US$ 3,8M", 0.25), ("Nova diária", "+78%", 0.45),
               ("Margem incr.", "→", 0.35), ("EBITDA novo", "US$ 15,3M", 0.55)],
              [GRAY, TEAL, ORANGE, NAVY])
    hero_callout(s, 8.35, EXHIBIT_Y + 0.4, 4.35, 2.9, "+303%",
                 "EBITDA do barco\nDiária +78%, custo fixo estável\nContratos até 2030", ORANGE)
    native_column(s, 0.75, EXHIBIT_Y + 3.45, 7.0, 0.95,
                  [("Rec.+8%", 8), ("EBITDA+30%", 30), ("Cob.88%", 88)],
                  [GRAY, TEAL, NAVY], ymax=100)

    # 14 Section 4 — Valuation
    section_slide(prs, 14, "4", "Valuation e catalisadores", "Múltiplo descontado, claims e fusão CBO")

    # 15 Valuation
    s = content_slide(prs, 15)
    action_lede(
        s,
        "Negociada a 3,7x EV/EBITDA, a OceanPact é a mais descontada do setor com caixa e claims visíveis",
        "Comparáveis de mercado — 26E",
        "Market Makers FIA",
        footnote="Ajuste exclui claims R$ 513M VP (~26% mkt cap)",
        exhibit_title="EV/EBITDA 26E — OceanPact vs referências (x)",
    )
    native_column(s, 0.75, EXHIBIT_Y + 0.45, 7.5, 3.0,
                  [("OPCT", 3.7), ("OPCT adj.", 3.1), ("WSUT deal", 4.4), ("Média setor", 7.5)],
                  [ORANGE, TEAL, GRAY, NAVY], fmt="0.0", ymax=9)
    hero_callout(s, 8.55, EXHIBIT_Y + 0.55, 3.85, 2.35, "3,7x",
                 "EV/EBITDA 26E\nMais descontada\nFCF yield elevado", NAVY)
    annotation_band(s, [
        "OPCT adj. 3,1x após claims — abaixo do deal WSUT (4,4x)",
        "Média setorial 7,5x — gap de repricing se ciclo se mantém",
    ], x=0.75, w=7.6)

    # 16 CBO
    s = content_slide(prs, 16)
    action_lede(
        s,
        "Fusão com a CBO cria o 2º maior OSV do Brasil sem diluir claims judiciais aos acionistas atuais",
        "Escala da frota — antes e depois",
        "OceanPact; Gradus",
        exhibit_title="Embarcações operacionais (unidades)",
    )
    native_column(s, 1.5, EXHIBIT_Y + 0.5, 5.0, 3.5,
                  [("OPCT", 28), ("CBO", 45), ("Combinada", 73)],
                  [TEAL, GRAY, NAVY])
    arrow_right(s, Inches(6.8), Inches(EXHIBIT_Y + 1.9), Inches(0.6), Inches(0.5), ORANGE)
    for i, (lbl, val, impl, col) in enumerate([
        ("Market share", "15%", "2º maior OSV do Brasil", NAVY),
        ("Sinergias/ano", "R$ 100M+", "Financia integração CBO", ORANGE),
        ("Acionistas OPCT", "42%", "UP Coral preservada", NAVY),
    ]):
        ty = EXHIBIT_Y + 0.45 + i * 1.3
        shape_rect(s, Inches(7.7), Inches(ty), Inches(5.0), Inches(1.15), WHITE, col, 2)
        text_at(s, 7.85, ty + 0.1, 4.7, 0.32, lbl, 12, False, GRAY)
        text_at(s, 7.85, ty + 0.42, 4.7, 0.42, val, 28, True, col)
        text_at(s, 7.85, ty + 0.88, 4.7, 0.28, impl, 10, False, GRAY)

    # 17 Upside
    s = content_slide(prs, 17)
    action_lede(
        s,
        "Só para igualar o múltiplo da WSUT, a OceanPact teria ~30% de upside; na média setorial, ~116%",
        "Cenários de repricing de múltiplo",
        "Tidewater; estimativas Market Makers",
        exhibit_title="Upside de repricing vs benchmarks de múltiplo (%)",
    )
    native_column(s, 2.0, EXHIBIT_Y + 0.45, 6.8, 3.2,
                  [("WSUT +30%", 30), ("Média +116%", 116)],
                  [TEAL, NAVY], ymax=130)
    hero_callout(s, 8.85, EXHIBIT_Y + 0.55, 3.55, 2.4, "+116%",
                 "Upside vs.\nmédia setorial 7,5x", NAVY)
    annotation_band(s, [
        "Tidewater pagou 4,4x na WSUT; OPCT negocia a 3,7x",
        "Disciplina: perdeu leilão RSV 4T24; reconquistou com diárias ~60% maiores",
    ])

    # 18 Risks + close
    s = content_slide(prs, 18)
    action_lede(
        s,
        "OceanPact oferece a melhor relação risco-retorno no offshore brasileiro — com disciplina em preço e execução CBO",
        "Riscos a monitorar e conclusão",
    )
    rows = [
        ("[Pilar 1] Capital intensity\nCapex ROV vs geração de caixa", ["Médio", "Alto"]),
        ("[Pilar 2] Petrobras 57%\nRenovação contratos e mix RSV", ["Alto", "Alto"]),
        ("[Pilar 3] Fusão CBO\nSinergias e alavancagem pós-deal", ["Alto", "Médio"]),
        ("[Macro] Ciclo óleo\nCapex Petrobras se Brent cair", ["Médio", "Médio"]),
        ("[Mercado] Overhang PE\nSaída Gradus e liquidez B3", ["Médio", "Médio"]),
    ]
    cell_colors = [
        [RGBColor(255, 242, 204), RGBColor(255, 199, 199)],
        [RGBColor(255, 199, 199), RGBColor(255, 150, 150)],
        [RGBColor(255, 199, 199), RGBColor(255, 242, 204)],
        [RGBColor(255, 242, 204), RGBColor(255, 242, 204)],
        [RGBColor(255, 242, 204), RGBColor(255, 242, 204)],
    ]
    matrix_risk(s, 0.65, EXHIBIT_Y + 0.25, 11.9, 3.35, rows, ["Impacto", "Prob."], cell_colors)
    shape_rect(s, Inches(0.65), Inches(5.15), Inches(11.9), Inches(1.15), WHITE, NAVY, 2)
    text_at(s, 0.85, 5.28, 11.5, 0.35, "Recomendação Market Makers FIA", 14, True, NAVY)
    text_at(s, 0.85, 5.65, 11.5, 0.55,
            "OceanPact é a porta de entrada preferencial no ciclo: escassez defensável, EBITDA alavancado, "
            "múltiplo descontado com claims e catalisador CBO.",
            12, False, GRAY)
    text_at(s, 0.55, 6.12, 12.2, 0.55,
            "Material educativo — Market Makers FIA (mai/2026). Não constitui recomendação de investimento.",
            8, False, GRAY)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved: {OUT} ({len(prs.slides)} slides, v6 storyline + native charts)")


if __name__ == "__main__":
    build()

