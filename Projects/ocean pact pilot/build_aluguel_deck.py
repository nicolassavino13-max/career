"""BUILD_DECK — Aluguel Moema/VNC executive deck (12 slides)."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(r"c:\Users\nicol\OneDrive\Área de Trabalho\NicOS\Projects\aluguel_moema_vnc_deck.pptx")

NAVY = RGBColor(0, 32, 96)
ORANGE = RGBColor(237, 125, 49)
GREEN = RGBColor(84, 130, 53)
RED = RGBColor(192, 0, 0)
GRAY = RGBColor(89, 89, 89)
LIGHT = RGBColor(242, 244, 248)
WHITE = RGBColor(255, 255, 255)
TEAL = RGBColor(5, 102, 152)


def set_run(run, size=18, bold=False, color=None, name="Arial"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, lines, accent):
    box = add_rect(slide, left, top, width, height, WHITE, accent)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(10)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run(r, 14, True, NAVY)
    for line in lines:
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        r = p.add_run()
        r.text = line
        set_run(r, 11, False, GRAY)


def add_title_bar(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.05), NAVY)
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.18), Inches(12), Inches(0.55))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run(r, 28, True, WHITE)
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.55), Inches(0.72), Inches(12), Inches(0.35))
        p = sb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        set_run(r, 14, False, RGBColor(200, 210, 230))


def add_bullets(slide, items, left=0.7, top=1.35, width=11.8, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(10 if i else 0)
        p.level = 0
        r = p.add_run()
        r.text = item
        set_run(r, size, False, GRAY)


def add_insight_row(slide, y, num, headline, detail, color):
    add_rect(slide, Inches(0.55), Inches(y), Inches(0.55), Inches(0.55), color)
    nb = slide.shapes.add_textbox(Inches(0.55), Inches(y + 0.08), Inches(0.55), Inches(0.4))
    p = nb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(num)
    set_run(r, 18, True, WHITE)
    hb = slide.shapes.add_textbox(Inches(1.25), Inches(y), Inches(11), Inches(0.35))
    r = hb.text_frame.paragraphs[0].add_run()
    r.text = headline
    set_run(r, 16, True, NAVY)
    db = slide.shapes.add_textbox(Inches(1.25), Inches(y + 0.38), Inches(11), Inches(0.45))
    r = db.text_frame.paragraphs[0].add_run()
    r.text = detail
    set_run(r, 12, False, GRAY)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 — Title
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), Inches(13.33), Inches(7.5), NAVY)
    add_rect(s, Inches(0), Inches(5.9), Inches(13.33), Inches(1.6), TEAL)
    t = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = "Morar perto do Ibirapuera"
    set_run(r, 40, True, WHITE)
    t2 = s.shapes.add_textbox(Inches(0.8), Inches(3.1), Inches(11.5), Inches(0.8))
    r = t2.text_frame.paragraphs[0].add_run()
    r.text = "Com o mínimo de ruído de Congonhas"
    set_run(r, 28, False, ORANGE)
    meta = s.shapes.add_textbox(Inches(0.8), Inches(6.15), Inches(11), Inches(0.5))
    r = meta.text_frame.paragraphs[0].add_run()
    r.text = "Busca Moema Pássaros + VNC miolo  ·  Maio 2026  ·  Orçamento ALL-IN ~R$ 7.500"
    set_run(r, 12, False, WHITE)

    # 2 — Storyline arc
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "A história em uma linha", "Objetivo → Onde → Achados → Visitas → Decisão")
    steps = [
        ("1", "Objetivo", "Parque + silêncio + orçamento"),
        ("2", "Onde", "Ruas certas > bairro certo"),
        ("3", "Achados", "3 portais, 10 imóveis ranqueados"),
        ("4", "Visitas", "Top 3 + teste de ruído"),
        ("5", "Decisão", "Confirmar vidro, andar, face do dorm"),
    ]
    x0 = 0.55
    for i, (num, label, sub) in enumerate(steps):
        x = x0 + i * 2.45
        add_rect(s, Inches(x), Inches(2.0), Inches(2.1), Inches(2.8), LIGHT, NAVY)
        circ = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.75), Inches(2.25), Inches(0.6), Inches(0.6))
        circ.fill.solid()
        circ.fill.fore_color.rgb = ORANGE if i == 2 else NAVY
        circ.line.fill.background()
        tb = s.shapes.add_textbox(Inches(x + 0.75), Inches(2.35), Inches(0.6), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        set_run(r, 16, True, WHITE)
        lb = s.shapes.add_textbox(Inches(x + 0.15), Inches(3.1), Inches(1.8), Inches(0.5))
        r = lb.text_frame.paragraphs[0].add_run()
        r.text = label
        set_run(r, 14, True, NAVY)
        sb = s.shapes.add_textbox(Inches(x + 0.15), Inches(3.55), Inches(1.8), Inches(1.0))
        r = sb.text_frame.paragraphs[0].add_run()
        r.text = sub
        set_run(r, 10, False, GRAY)
    note = s.shapes.add_textbox(Inches(0.55), Inches(5.2), Inches(12), Inches(0.6))
    r = note.text_frame.paragraphs[0].add_run()
    r.text = "Mensagem-chave: ruído de avião é filtro #1 — preço e m² vêm depois."
    set_run(r, 14, True, TEAL)

    # 3 — Objective
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "O que estamos otimizando", "Três critérios — nesta ordem")
    cards = [
        ("1  Ruído", ["Miolo de quadra", "Ruas prioritárias", "Evitar avenidas / Índios"]),
        ("2  Isolamento", ["Vidro laminado/duplo", "Prédio 2014+ como proxy", "Dorm para o miolo"]),
        ("3  Vida no bairro", ["~10 min Ibirapuera a pé", "1 dorm · 35–60 m²", "ALL-IN ≤ R$ 7.500"]),
    ]
    for i, (title, lines) in enumerate(cards):
        add_card(s, Inches(0.55 + i * 4.15), Inches(1.5), Inches(3.85), Inches(4.5), title, lines, ORANGE if i == 0 else NAVY)

    # 4 — Where to focus
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Onde focar no mapa", "Não é só o bairro — é a rua e a face do prédio")
    add_rect(s, Inches(0.55), Inches(1.45), Inches(5.8), Inches(4.8), RGBColor(232, 245, 233))
    gt = s.shapes.add_textbox(Inches(0.8), Inches(1.65), Inches(5.2), Inches(0.4))
    r = gt.text_frame.paragraphs[0].add_run()
    r.text = "PRIORIZAR"
    set_run(r, 16, True, GREEN)
    add_bullets(s, [
        "Moema Pássaros — oeste da Av. Ibirapuera",
        "VNC miolo — Afonso Braz, Domingos Leme",
        "Ruas: Canário, Rouxinol, Pintassilgo, Pavão",
        "Dormitório para pátio / miolo de quadra",
    ], left=0.8, top=2.1, width=5.2, size=14)
    add_rect(s, Inches(6.9), Inches(1.45), Inches(5.8), Inches(4.8), RGBColor(255, 235, 235))
    rt = s.shapes.add_textbox(Inches(7.15), Inches(1.65), Inches(5.2), Inches(0.4))
    r = rt.text_frame.paragraphs[0].add_run()
    r.text = "EVITAR"
    set_run(r, 16, True, RED)
    add_bullets(s, [
        "Moema Índios — corredor norte de aviões",
        "Frente Av. Moema, Ibirapuera, Santo Amaro",
        "Andar baixo em zona exposta",
        "Diogo Jácome perto Bandeirantes",
    ], left=7.15, top=2.1, width=5.2, size=14)
    tip = s.shapes.add_textbox(Inches(0.55), Inches(6.35), Inches(12), Inches(0.5))
    r = tip.text_frame.paragraphs[0].add_run()
    r.text = "Congonhas aproxima pelo norte → miolo oeste e quadra densa amortecem o corredor."
    set_run(r, 12, True, GRAY)

    # 5 — How we searched
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Como buscamos", "Filtros estruturados + busca por ruas prioritárias")
    portals = [
        ("QuintoAndar", "Melhor encaixe\nrua + orçamento", GREEN),
        ("ZAP Imóveis", "Destaque em\nPintassilgo", ORANGE),
        ("Imovelweb", "Sem listagens\nauditáveis", RED),
    ]
    for i, (name, desc, col) in enumerate(portals):
        x = 0.55 + i * 4.15
        add_rect(s, Inches(x), Inches(1.8), Inches(3.85), Inches(3.2), WHITE, col)
        add_rect(s, Inches(x), Inches(1.8), Inches(3.85), Inches(0.55), col)
        nb = s.shapes.add_textbox(Inches(x + 0.2), Inches(1.9), Inches(3.4), Inches(0.4))
        r = nb.text_frame.paragraphs[0].add_run()
        r.text = name
        set_run(r, 16, True, WHITE)
        db = s.shapes.add_textbox(Inches(x + 0.2), Inches(2.6), Inches(3.4), Inches(2.0))
        r = db.text_frame.paragraphs[0].add_run()
        r.text = desc
        set_run(r, 20, True, NAVY)
    add_bullets(s, [
        "ALL-IN = Aluguel + Condomínio + IPTU (sem taxa portal)",
        "Vidro acústico: busca por linguagem natural → zero hits explícitos",
        "10 imóveis ranqueados · Top 3 prontos para visita",
    ], left=0.55, top=5.3, width=12, size=13)

    # 6 — Three insights
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Três insights da busca")
    add_insight_row(s, 1.45, 1, "Oferta ideal concentra-se no QuintoAndar",
                    "Moema Pássaros e VNC miolo — melhor relação rua prioritária + ALL-IN.", GREEN)
    add_insight_row(s, 2.55, 2, "Nenhum anúncio cita vidro acústico",
                    "Campo quase sempre \"não informado\". Prédio 2014+ (ex. Today Moema) é proxy, não prova.", ORANGE)
    add_insight_row(s, 3.65, 3, "ZAP traz o melhor custo em rua prioritária",
                    "Pintassilgo 50 m² — ALL-IN R$ 3.612. Imovelweb precisa rodada manual.", TEAL)

    # 7 — Top 3 visits
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Top 3 para visitar primeiro", "Shortlist — ruído baixo + rua prioritária + dentro do teto")
    top3 = [
        ("#1  Rua Canário · 37 m²", "ALL-IN ~R$ 3.422", "Risco baixo · miolo Pássaros", GREEN),
        ("#2  Av. Rouxinol · 43 m²", "ALL-IN ~R$ 5.313", "12º–15º · reformado · mobiliado", NAVY),
        ("#3  Rua Pintassilgo · 50 m²", "ALL-IN R$ 3.612", "ZAP · melhor custo/rua", ORANGE),
    ]
    for i, (title, price, note, col) in enumerate(top3):
        x = 0.55 + i * 4.15
        add_rect(s, Inches(x), Inches(1.55), Inches(3.85), Inches(4.6), WHITE, col)
        add_rect(s, Inches(x), Inches(1.55), Inches(3.85), Inches(0.12), col)
        tb = s.shapes.add_textbox(Inches(x + 0.2), Inches(1.85), Inches(3.4), Inches(1.0))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = title
        set_run(r, 15, True, NAVY)
        pb = s.shapes.add_textbox(Inches(x + 0.2), Inches(2.9), Inches(3.4), Inches(0.8))
        r = pb.text_frame.paragraphs[0].add_run()
        r.text = price
        set_run(r, 26, True, col)
        nb = s.shapes.add_textbox(Inches(x + 0.2), Inches(3.8), Inches(3.4), Inches(1.5))
        r = nb.text_frame.paragraphs[0].add_run()
        r.text = note
        set_run(r, 12, False, GRAY)

    # 8 — Risk lens
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Lente de risco", "Preço baixo não compensa exposição ao corredor aéreo")
    matrix = [
        ("Baixo risco", "Canário, Pintassilgo, Rouxinol (alto)", GREEN),
        ("Médio", "Today Moema / Pavão — prédio novo, via aberta", ORANGE),
        ("Alto — descartar salvo prova", "Moema Índios, Av. Moema 3º, Av. Iraí", RED),
    ]
    for i, (label, examples, col) in enumerate(matrix):
        y = 1.55 + i * 1.55
        add_rect(s, Inches(0.55), Inches(y), Inches(2.2), Inches(1.2), col)
        lb = s.shapes.add_textbox(Inches(0.65), Inches(y + 0.35), Inches(2.0), Inches(0.5))
        r = lb.text_frame.paragraphs[0].add_run()
        r.text = label
        set_run(r, 12, True, WHITE)
        ex = s.shapes.add_textbox(Inches(3.0), Inches(y + 0.2), Inches(9.5), Inches(0.9))
        r = ex.text_frame.paragraphs[0].add_run()
        r.text = examples
        set_run(r, 16, False, GRAY)

    # 9 — Broker questions
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "4 perguntas para o corretor", "Confirmar antes de reservar visita")
    qs = [
        "As janelas do dorm são laminado ou duplo? Qual espessura?",
        "Em que ano o prédio foi entregue?",
        "O dorm abre para miolo de quadra ou para rua/avenida?",
        "Qual andar exato da unidade?",
    ]
    for i, q in enumerate(qs):
        y = 1.55 + i * 1.15
        add_rect(s, Inches(0.55), Inches(y), Inches(0.45), Inches(0.45), ORANGE)
        nb = s.shapes.add_textbox(Inches(0.55), Inches(y + 0.06), Inches(0.45), Inches(0.35))
        p = nb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(i + 1)
        set_run(r, 14, True, WHITE)
        qb = s.shapes.add_textbox(Inches(1.15), Inches(y + 0.05), Inches(11.5), Inches(0.5))
        r = qb.text_frame.paragraphs[0].add_run()
        r.text = q
        set_run(r, 16, False, GRAY)
    quote = s.shapes.add_textbox(Inches(0.55), Inches(6.0), Inches(12), Inches(0.8))
    r = quote.text_frame.paragraphs[0].add_run()
    r.text = '"Posso visitar nos horários de pico de voos, com janelas fechadas e abertas?"'
    set_run(r, 14, True, NAVY)

    # 10 — Visit protocol
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Protocolo de visita", "O teste de ruído é inegociável")
    steps_v = [
        ("Horário", "Manhã cedo + fim de tarde, dias úteis (pico Congonhas)"),
        ("No quarto", "Janelas fechadas → ouvir 3 min → abrir → comparar"),
        ("Ferramenta", "Decibelímetro no celular — registrar pico"),
        ("Descarte", "Se dorm exposto + andar baixo em zona Índios → sair"),
    ]
    for i, (label, detail) in enumerate(steps_v):
        y = 1.6 + i * 1.2
        add_rect(s, Inches(0.55), Inches(y), Inches(1.6), Inches(0.85), LIGHT, NAVY)
        lb = s.shapes.add_textbox(Inches(0.65), Inches(y + 0.22), Inches(1.4), Inches(0.4))
        r = lb.text_frame.paragraphs[0].add_run()
        r.text = label
        set_run(r, 12, True, NAVY)
        db = s.shapes.add_textbox(Inches(2.35), Inches(y + 0.15), Inches(10), Inches(0.6))
        r = db.text_frame.paragraphs[0].add_run()
        r.text = detail
        set_run(r, 15, False, GRAY)

    # 11 — Next steps
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Próximos passos", "Da busca à decisão em 5 movimentos")
    actions = [
        "Fechar orçamento ALL-IN real e preferência mobiliado",
        "Agendar visitas #1–3 nos horários de pico",
        "Repetir Imovelweb manualmente (VNC + ruas prioritárias)",
        "Criar alertas QuintoAndar/ZAP nas ruas da shortlist",
        "Descartar Índios e avenidas salvo teste comprovar ruído baixo",
    ]
    for i, act in enumerate(actions):
        y = 1.55 + i * 0.95
        add_rect(s, Inches(0.55), Inches(y), Inches(12.2), Inches(0.75), WHITE if i % 2 else LIGHT)
        nb = s.shapes.add_textbox(Inches(0.7), Inches(y + 0.15), Inches(0.5), Inches(0.4))
        r = nb.text_frame.paragraphs[0].add_run()
        r.text = str(i + 1)
        set_run(r, 16, True, ORANGE)
        ab = s.shapes.add_textbox(Inches(1.25), Inches(y + 0.15), Inches(11), Inches(0.5))
        r = ab.text_frame.paragraphs[0].add_run()
        r.text = act
        set_run(r, 15, False, GRAY)

    # 12 — Decision frame
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Como decidir", "Framework simples pós-visita")
    add_rect(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(4.8), NAVY)
    lines = [
        "Ruído no dorm (pico)  →  peso 50%",
        "Vidro / ano do prédio  →  peso 25%",
        "ALL-IN + mobiliado  →  peso 15%",
        "Distância Ibirapuera / metrô  →  peso 10%",
    ]
    for i, line in enumerate(lines):
        tb = s.shapes.add_textbox(Inches(1.0), Inches(2.0 + i * 0.85), Inches(11), Inches(0.6))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = line
        set_run(r, 20, i == 0, WHITE)
    close = s.shapes.add_textbox(Inches(0.55), Inches(6.35), Inches(12), Inches(0.5))
    r = close.text_frame.paragraphs[0].add_run()
    r.text = "Recomendação: visitar Canário, Rouxinol e Pintassilgo esta semana."
    set_run(r, 14, True, TEAL)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
